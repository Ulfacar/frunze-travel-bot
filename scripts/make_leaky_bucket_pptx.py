"""Сборка PPTX презентации «Дырявое ведро» (Frunze Travel).

Повторяет структуру HTML-деки v4 (порядок «боль до карты» + честные
формулировки). Переиспользует движок стилей из make_presentation.py.

Запуск:  python scripts/make_leaky_bucket_pptx.py
Результат: C:\\Users\\alanb\\OneDrive\\Рабочий стол\\Frunze-дырявое-ведро.pptx
"""
from __future__ import annotations

import os
import sys
from pathlib import Path

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from make_presentation import (  # noqa: E402
    Palette, SAFE_L, SAFE_R, add_bg, add_card, add_shape, add_text,
)

OUT = Path(r"C:\Users\alanb\OneDrive\Рабочий стол\Frunze-дырявое-ведро.pptx")
FOOT = "Frunze Travel · GetVisa · Бишкек"


def _blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    return s


def _foot(slide, n, dark=False):
    c = Palette.muted_dark if dark else Palette.slate
    add_text(slide, SAFE_L, 6.66, 7.0, 0.25, FOOT, 8.5, False, c)
    add_text(slide, 11.7, 6.66, 1.0, 0.25, f"{n:02d}", 8.5, True, c, PP_ALIGN.RIGHT)


def _head(slide, kicker, title, dark=False):
    kc = Palette.muted_dark if dark else Palette.teal_dark
    tc = Palette.white if dark else Palette.ink
    add_text(slide, SAFE_L, 0.52, 11.0, 0.28, kicker.upper(), 10.5, True, kc)
    add_text(slide, SAFE_L, 0.86, 11.5, 1.0, title, 27, True, tc)
    add_shape(slide, MSO_SHAPE.RECTANGLE, SAFE_L, 1.72, 2.2, 0.055, Palette.teal, label="rule")


def _bubble(slide, x, y, w, h, text, fill, color, italic=False, size=11.5):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, label="bubble")
    add_text(slide, x + 0.12, y + 0.06, w - 0.24, h - 0.12, text, size, False, color,
             italic=italic, valign=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------- slides

def s1_cover(prs):
    s = _blank(prs); add_bg(s, dark=True)
    add_shape(s, MSO_SHAPE.RECTANGLE, SAFE_L, 0.9, 2.2, 0.06, Palette.amber, label="a")
    add_text(s, SAFE_L, 1.15, 11.0, 0.3, "FRUNZE TRAVEL · GETVISA · РАЗБОР ВОРОНКИ", 11, True, Palette.amber)
    add_text(s, SAFE_L, 1.7, 11.6, 2.4,
             "Проблема не в потоке.\nПроблема в том, что вы\nне видите, где он утекает.",
             38, True, Palette.white)
    add_text(s, SAFE_L, 4.35, 10.8, 1.1,
             "Реклама льёт, бот отвечает, клиенты пишут. Но между «написал» и «оплатил» "
             "стоит ночь, очередь и уставший менеджер — и деньги уходят там, где никто не смотрит.",
             15, False, Palette.muted_dark)
    for i, (t) in enumerate(["Не больше лидов — меньше потерь", "Контроль над воронкой",
                             "×2 на том же бюджете"]):
        x = SAFE_L + i * 4.0
        add_shape(s, MSO_SHAPE.OVAL, x, 5.78, 0.12, 0.12, Palette.teal, label="d")
        add_text(s, x + 0.22, 5.66, 3.7, 0.4, t, 12, True, Palette.white)
    _foot(s, 1, dark=True)


def s2_honesty(prs):
    s = _blank(prs); add_bg(s)
    _head(s, "Честно · ваши реальные цифры", "Мы уже измерили. Вот что показала ваша база.")
    # concrete (dark card)
    add_card(s, SAFE_L, 1.95, 5.75, 3.7, Palette.navy, Palette.navy)
    add_text(s, SAFE_L + 0.3, 2.12, 5.2, 0.3, "БЕТОН — ИЗМЕРЕНО НА ВАШЕЙ БАЗЕ", 11, True, Palette.amber)
    for i, t in enumerate([
        "Медиана ответа менеджера — 3 часа",
        "66 лидов/день (медиана; пик 96), не «сто»",
        "20% лидов рождаются ночью (22:00–08:00)",
        "157 лидов застоялись без ответа >2 суток",
        "74% диалогов бот ведёт сам, без человека",
    ]):
        y = 2.55 + i * 0.6
        add_text(s, SAFE_L + 0.3, y, 0.3, 0.4, "✓", 13, True, Palette.amber)
        add_text(s, SAFE_L + 0.62, y, 4.9, 0.55, t, 11.5, False, Palette.white)
    # hypothesis (light card)
    hx = 6.6
    add_card(s, hx, 1.95, 5.75, 3.7, Palette.white, Palette.teal)
    add_text(s, hx + 0.3, 2.12, 5.2, 0.3, "СЛЕПАЯ ЗОНА — ВЫ ЭТО НЕ ВИДИТЕ", 11, True, Palette.teal_dark)
    for i, t in enumerate([
        "Конверсия: размечено 3 из 389 — вы её не считаете",
        "Источник рекламы захвачен на 0% — что сработало?",
        "16% лидов «готовы платить» — лежат необработанными",
        "Сколько из 157 застойных ещё живы",
        "Реальная выручка с лида — нет источника истины",
    ]):
        y = 2.55 + i * 0.6
        add_text(s, hx + 0.3, y, 0.3, 0.4, "?", 13, True, Palette.teal)
        add_text(s, hx + 0.62, y, 4.9, 0.55, t, 11.5, False, Palette.ink)
    # disarm
    add_card(s, SAFE_L, 5.85, 12.13, 0.72, Palette.amber_soft, Palette.amber)
    add_text(s, SAFE_L + 0.25, 5.95, 11.6, 0.55,
             "Раньше тут висели цифры рекламщика ($3.8/лид, $100/тур, 100 лидов) — мы заменили их ВАШИМИ, "
             "снятыми с прода 06.07. Главное: из 389 диалогов размечено 3 продажи — не потому что не продаёте, "
             "а потому что конверсию сейчас никто не считает. «Вы не видите, где течёт» — это буквально.",
             10.5, False, Palette.ink, valign=MSO_ANCHOR.MIDDLE)
    _foot(s, 2)


def s3_night(prs):
    s = _blank(prs); add_bg(s)
    _head(s, "Главная дыра · одна, ваша", "Ночь: лучший спрос — там, где нет света.")
    # chat log (left)
    _bubble(s, SAFE_L, 2.0, 5.4, 0.72, "Здравствуйте! Хотим на море, вдвоём, в сентябре. Какие варианты?  · 23:47",
            Palette.grey_soft, Palette.ink)
    _bubble(s, SAFE_L + 1.3, 2.85, 4.5, 0.5, "— тишина: менеджер спит, бот отдал хендофф —",
            Palette.red_soft, Palette.red, italic=True, size=11)
    _bubble(s, SAFE_L, 3.5, 5.0, 0.55, "Ну что там? Мы бы хотели определиться сегодня  · 00:12",
            Palette.grey_soft, Palette.ink)
    _bubble(s, SAFE_L + 1.0, 4.2, 4.8, 0.72, "Добрый день! Подскажите даты и бюджет 🙂  · СЛЕДУЮЩИЙ вечер",
            Palette.navy, Palette.white)
    add_text(s, SAFE_L, 5.05, 5.6, 0.4, "…прочитано. Ответа нет — купили в другом месте.",
             11, True, Palette.red, italic=True)
    # right column
    rx = 6.9
    add_text(s, rx, 2.05, 5.6, 0.5, "Это в ваших данных", 18, True, Palette.ink)
    add_text(s, rx, 2.7, 5.6, 1.6,
             "20% ваших лидов рождаются ночью (78 из 389). А медиана ответа менеджера — 3 часа. Пока "
             "человек спит или занят, тёплый ночной клиент остывает и уходит: самый горячий момент "
             "приходится на время, когда за экраном никого нет.", 13, False, Palette.slate)
    add_card(s, rx, 4.5, 5.6, 1.35, Palette.green_soft, Palette.teal)
    add_text(s, rx + 0.25, 4.65, 5.1, 1.1,
             "Проколоть эту дыру = поставить того, кто никогда не спит: ловит клиента на пике и держит "
             "до утра, когда придёт человек. Это ядро решения.", 12.5, False, Palette.teal_dark,
             valign=MSO_ANCHOR.MIDDLE)
    _foot(s, 3)


def s4_more(prs):
    s = _blank(prs); add_bg(s)
    _head(s, "Принцип один · дыр несколько", "…и таких точек ещё восемь.")
    add_text(s, SAFE_L, 1.85, 11.8, 0.4,
             "Ночь — самая крупная. Принцип везде один: клиент уже пришёл и почти готов — надо не потерять.",
             13, False, Palette.slate)
    items = [
        ("1", "Три касания", "Нет системы «догнать» — лид остывает."),
        ("2", "Консультация ≠ продажа", "Всё рассказали — цену унесли к конкуренту."),
        ("3", "157 застойных", "Трафик лежит без касания дольше 2 суток."),
        ("4", "Билеты жрут часы", "2–3 часа тур-менеджера на нерешительного."),
        ("5", "Голос — в тишину", "Войс = интерес, а в ответ хендофф и пауза."),
        ("6", "Нет слотов", "Никто не звонит вовремя — время не назначали."),
        ("7", "Бюджет размазан", "Поровну на неравные воронки."),
        ("8", "Две панели", "Своя + Bitrix = двойная работа на фиче."),
    ]
    cw, ch, gx, gy = 2.86, 1.55, 0.23, 0.25
    x0, y0 = SAFE_L, 2.45
    for i, (n, h, d) in enumerate(items):
        col, row = i % 4, i // 4
        x = x0 + col * (cw + gx)
        y = y0 + row * (ch + gy)
        add_card(s, x, y, cw, ch, Palette.white, Palette.line)
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.2, y + 0.2, 0.4, 0.4, Palette.green_soft, label="b")
        add_text(s, x + 0.2, y + 0.23, 0.4, 0.36, n, 13, True, Palette.teal_dark, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.2, y + 0.68, cw - 0.4, 0.34, h, 12.5, True, Palette.ink)
        add_text(s, x + 0.2, y + 1.02, cw - 0.4, 0.45, d, 10, False, Palette.slate)
    _foot(s, 4)


def s5_model(prs):
    s = _blank(prs); add_bg(s)
    _head(s, "Как делаются деньги", "Три воронки — у каждой свой характер.")
    cols = [
        ("Визы", "чистые деньги", Palette.teal, [
            ("USA-виза", "самые чистые деньги"), ("Путь", "лид → квал → оплата"),
            ("Риск", "растёт конкуренция")]),
        ("Туры", "крупные деньги, ночью", Palette.amber, [
            ("Где", "ночь, голос, человек"), ("Квал", "деньги + даты + ЛПР"),
            ("Заработок", "на возврате клиента")]),
        ("Билеты", "стабильный день", Palette.teal_dark, [
            ("Деньги", "малые, но каждый день"), ("Решение", "выделить 1 менеджера"),
            ("Даёт", "часы под туры")]),
    ]
    cw = 3.8
    for i, (name, tag, accent, rows) in enumerate(cols):
        x = SAFE_L + i * (cw + 0.36)
        add_card(s, x, 2.0, cw, 3.15, Palette.white, Palette.line)
        add_shape(s, MSO_SHAPE.RECTANGLE, x, 2.0, cw, 0.08, accent, label="top")
        add_text(s, x + 0.25, 2.25, cw - 0.5, 0.4, name, 18, True, Palette.ink)
        add_text(s, x + 0.25, 2.68, cw - 0.5, 0.3, tag.upper(), 9.5, True, accent)
        for j, (k, v) in enumerate(rows):
            y = 3.25 + j * 0.6
            add_text(s, x + 0.25, y, 1.5, 0.5, k, 10.5, False, Palette.slate)
            add_text(s, x + 1.6, y, cw - 1.8, 0.5, v, 11, True, Palette.ink, align=PP_ALIGN.RIGHT)
    add_card(s, SAFE_L, 5.4, 12.13, 0.95, Palette.grey_soft, Palette.line)
    add_text(s, SAFE_L + 0.25, 5.52, 11.6, 0.75,
             "Важно про билеты: их не «убивают» — они приносят стабильные деньги каждый день. "
             "Задача — снять их с плеч тур-менеджеров: один человек на билеты, двое освобождаются под крупные туры.",
             12, False, Palette.ink, valign=MSO_ANCHOR.MIDDLE)
    _foot(s, 5)


def s6_risk(prs):
    s = _blank(prs); add_bg(s, dark=True)
    add_text(s, SAFE_L, 0.55, 11.0, 0.28, "ЧЕСТНЫЙ РИСК · ДЫРА В ДНЕ", 10.5, True, Palette.amber)
    add_text(s, SAFE_L, 0.9, 11.8, 0.7, "А менеджеры вообще будут это делать?", 27, True, Palette.white)
    add_text(s, SAFE_L, 1.9, 11.8, 0.9,
             "«Если менеджер игнорит живого лида сегодня — почему не проигнорит слот завтра в 9 утра?»",
             17, True, Palette.amber, italic=True)
    add_text(s, SAFE_L, 2.95, 11.8, 0.6,
             "Справедливо. Поэтому решение не про «замотивируем», а про структуру: деньги ловятся даже "
             "при слабом менеджере, а слабый — становится виден.", 13, False, Palette.muted_dark)
    steps = [
        ("1", "Слот с окном", "Не взял звонок в течение ~получаса — слот не растворяется, а поднимает флаг."),
        ("2", "Эскалация", "Просроченный слот уходит другому менеджеру, вам — алерт в телефон."),
        ("3", "Табло утечки", "Поимённо: выдано / отработано / сгорело. Вы впервые видите, кто течёт."),
    ]
    cw = 3.8
    for i, (n, t, d) in enumerate(steps):
        x = SAFE_L + i * (cw + 0.36)
        add_card(s, x, 3.85, cw, 2.4, Palette.panel, Palette.panel)
        add_text(s, x + 0.28, 4.05, 1.0, 0.6, n, 26, True, Palette.amber)
        add_text(s, x + 0.28, 4.75, cw - 0.5, 0.4, t, 15, True, Palette.white)
        add_text(s, x + 0.28, 5.2, cw - 0.55, 0.95, d, 11.5, False, Palette.muted_dark)
    _foot(s, 6, dark=True)


def s7_dispatcher(prs):
    s = _blank(prs); add_bg(s)
    _head(s, "Ядро решения", "Ночной диспетчер: ловит ночью, вооружает утром.")
    steps = [
        ("Клиент пишет в 23:47", "Бот отвечает мгновенно — пока менеджер спит."),
        ("Квалификация по 4 критериям", "Деньги · даты · потребность · ЛПР (кто платит)."),
        ("Бронь слота на утро", "«Завтра в 13:00» — с окном и эскалацией."),
        ("К 9:00 — ДОСЬЕ менеджеру", "Бюджет, семья, куда мечтает, срочность. Войс расшифрован."),
    ]
    for i, (t, d) in enumerate(steps):
        y = 2.05 + i * 1.05
        accent = Palette.amber if i == 3 else Palette.teal
        add_shape(s, MSO_SHAPE.OVAL, SAFE_L, y, 0.5, 0.5, accent, label="n")
        add_text(s, SAFE_L, y + 0.06, 0.5, 0.4, str(i + 1), 15, True, Palette.white, align=PP_ALIGN.CENTER)
        add_text(s, SAFE_L + 0.7, y - 0.02, 5.4, 0.4, t, 14, True, Palette.ink)
        add_text(s, SAFE_L + 0.7, y + 0.38, 5.4, 0.55, d, 11, False, Palette.slate)
    # KPI cards (right)
    kx = 6.9
    kpis = [
        ("KPI БОТА", "% слотов, где досье попало в цель", Palette.teal, Palette.green_soft),
        ("KPI МЕНЕДЖЕРА", "% слотов, ставших деньгами", Palette.amber, Palette.amber_soft),
        ("БОНУС", "Прозрачность стоимости лида по воронкам", Palette.teal_dark, Palette.grey_soft),
    ]
    for i, (who, m, ac, bg) in enumerate(kpis):
        y = 2.0 + i * 1.5
        add_card(s, kx, y, 5.6, 1.3, bg, ac)
        add_text(s, kx + 0.28, y + 0.15, 5.0, 0.3, who, 10, True, ac)
        add_text(s, kx + 0.28, y + 0.5, 5.0, 0.65, m, 14, True, Palette.ink, valign=MSO_ANCHOR.MIDDLE)
    _foot(s, 7)


def s8_second(prs):
    s = _blank(prs); add_bg(s)
    _head(s, "Горизонт 2 · второй ×2, дешевле первого", "Вы уже купили этих клиентов — ценой ниже нуля.")
    add_card(s, SAFE_L, 2.2, 12.13, 3.0, Palette.amber_soft, Palette.amber)
    add_text(s, SAFE_L + 0.4, 2.5, 2.6, 2.4, "×2", 90, True, Palette.amber, valign=MSO_ANCHOR.MIDDLE)
    add_text(s, SAFE_L + 3.3, 2.55, 8.4, 1.5,
             "Ваша модель — не разовая продажа, а захват клиента: даёте цену ниже, чтобы он остался вашим. "
             "Первый ×2 — ночной диспетчер на входящем потоке. Второй ×2 — машина возврата: та же база, "
             "разбуженная под сезон и горящий тур. Трафик за неё уже оплачен минусовой ценой.",
             14, False, Palette.ink)
    add_text(s, SAFE_L + 3.3, 4.15, 8.4, 0.9,
             "«Если бы вам запретили купить хоть одного нового лида в этом году — что бы вы сделали "
             "с теми, кого уже купили в минус?»", 15, True, Palette.teal_dark, italic=True)
    add_text(s, SAFE_L, 5.5, 12.0, 0.7,
             "Осторожно: возврат делаем не авто-спамом с боевого номера (бан WhatsApp = весь бизнес разом), "
             "а как задачи менеджерам + аккуратные касания. Безопасность номера важнее скорости.",
             11, False, Palette.slate, italic=True)
    _foot(s, 8)


def s9_funnel(prs):
    s = _blank(prs); add_bg(s)
    _head(s, "Модель · пока на вилках", "Как это выглядит в воронке.")
    add_text(s, SAFE_L, 1.8, 11.8, 0.5,
             "Числа ниже — иллюстрация механики, не ваш замер (меряем с 1-й ночи). Показываем не точную "
             "цифру, а рычаг: где теряется и где можно удержать.", 12.5, False, Palette.slate)

    def col(x, title, accent, rows, tag):
        add_card(s, x, 2.55, 5.6, 3.35, Palette.white, Palette.line)
        add_shape(s, MSO_SHAPE.RECTANGLE, x, 2.55, 5.6, 0.08, accent, label="t")
        add_text(s, x + 0.28, 2.75, 4.0, 0.4, title, 15, True, Palette.ink)
        add_text(s, x + 3.8, 2.78, 1.6, 0.35, tag.upper(), 9, True, accent, align=PP_ALIGN.RIGHT)
        for j, (lbl, pct) in enumerate(rows):
            y = 3.35 + j * 0.5
            add_text(s, x + 0.28, y - 0.02, 3.6, 0.3, lbl, 10, False, Palette.slate)
            add_text(s, x + 4.7, y - 0.02, 0.7, 0.3, f"{pct}%", 10, True, Palette.ink, align=PP_ALIGN.RIGHT)
            add_shape(s, MSO_SHAPE.RECTANGLE, x + 0.28, y + 0.24, 5.05, 0.14, Palette.grey_soft, label="bg")
            add_shape(s, MSO_SHAPE.RECTANGLE, x + 0.28, y + 0.24, max(0.05, 5.05 * pct / 100), 0.14, accent, label="fg")

    col(SAFE_L, "Сегодня", Palette.teal,
        [("Входящие", 100), ("Внятный диалог", 60), ("Квал", 20), ("Звонок/встреча", 8), ("Продажа", 5)], "как есть")
    col(6.93, "Достижимо", Palette.amber,
        [("Ночь поймана", 85), ("Квал по 4 критериям", 30), ("Бронь слота", 22), ("Закрыл с досье", 10), ("Продажа", 10)], "потенциал")
    add_text(s, SAFE_L, 6.05, 12.0, 0.5,
             "Медиана ответа сегодня — 3 часа; ночью бот отвечает за секунды. + 157 застойных — резерв сверху. "
             "Точную дельту в $ поставим, как только начнём размечать продажи (сейчас 3 из 389).", 11, True, Palette.teal_dark)
    _foot(s, 9)


def s10_monday(prs):
    s = _blank(prs); add_bg(s, dark=True)
    add_text(s, SAFE_L, 0.55, 11.0, 0.28, "ЧТО ДЕЛАЕМ ДАЛЬШЕ", 10.5, True, Palette.amber)
    add_text(s, SAFE_L, 0.9, 11.8, 0.7, "Сначала свет и счётчик, потом ночная смена.", 26, True, Palette.white)
    left = [
        ("Решаете вы", "Цена туров: вилка или точная?"),
        ("Решаете вы", "Одна главная панель: Bitrix или своя?"),
        ("Решаете вы", "Билеты: выделяем менеджера?"),
    ]
    right = [
        ("Строим мы", "Потушить прод + включить счётчик $/лид"),
        ("Строим мы", "Ночной диспетчер: квал → слот → досье + табло"),
        ("Горизонт 2", "Машина возврата — разбудить оплаченную базу"),
    ]
    for i, (tag, t) in enumerate(left):
        y = 2.0 + i * 0.95
        add_card(s, SAFE_L, y, 5.75, 0.8, Palette.panel, Palette.panel)
        add_text(s, SAFE_L + 0.25, y + 0.1, 5.2, 0.25, tag.upper(), 8.5, True, Palette.amber)
        add_text(s, SAFE_L + 0.25, y + 0.35, 5.2, 0.4, t, 12.5, True, Palette.white)
    for i, (tag, t) in enumerate(right):
        y = 2.0 + i * 0.95
        add_card(s, 6.6, y, 5.75, 0.8, Palette.panel, Palette.panel)
        add_text(s, 6.85, y + 0.1, 5.2, 0.25, tag.upper(), 8.5, True, Palette.teal)
        add_text(s, 6.85, y + 0.35, 5.2, 0.4, t, 12.5, True, Palette.white)
    add_text(s, SAFE_L, 5.15, 12.0, 1.2,
             "«Мне не нужно больше лидов. Мне нужно перестать терять тех, за кого я уже заплатил — "
             "и видеть, где именно течёт.»", 18, True, Palette.white, italic=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)
    _foot(s, 10, dark=True)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for fn in (s1_cover, s2_honesty, s3_night, s4_more, s5_model, s6_risk,
               s7_dispatcher, s8_second, s9_funnel, s10_monday):
        fn(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"OK: {len(prs.slides._sldIdLst)} слайдов → {OUT}")


if __name__ == "__main__":
    build()
