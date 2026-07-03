from pathlib import Path
import base64
import shutil
import subprocess
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")


OUT_PATH = Path(r"C:\Users\alanb\OneDrive\Рабочий стол\Frunze-презентация.pptx")
FINAL_PATH = Path(r"C:\Users\alanb\OneDrive\Рабочий стол\Frunze-презентация-final.pptx")
BUILD_PATH = Path(__file__).resolve().parents[1] / "frunze_presentation_build.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5
SAFE_L = 0.6
SAFE_T = 0.5
SAFE_R = 12.73
SAFE_B = 7.0
CONTENT_W = SAFE_R - SAFE_L
FONT = "Arial"


class Palette:
    navy = RGBColor(12, 31, 46)
    panel = RGBColor(18, 45, 66)
    ink = RGBColor(21, 36, 61)
    slate = RGBColor(84, 99, 124)
    teal = RGBColor(0, 156, 148)
    teal_dark = RGBColor(0, 103, 109)
    green_soft = RGBColor(224, 245, 241)
    amber = RGBColor(245, 180, 64)
    amber_soft = RGBColor(255, 244, 218)
    red = RGBColor(222, 74, 77)
    red_soft = RGBColor(253, 232, 232)
    grey = RGBColor(138, 149, 163)
    grey_soft = RGBColor(238, 242, 246)
    paper = RGBColor(247, 249, 251)
    white = RGBColor(255, 255, 255)
    line = RGBColor(219, 226, 235)
    muted_dark = RGBColor(183, 204, 211)


def emu(value):
    return Inches(value)


def rgb_hex(color):
    return f"#{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def ensure_box(left, top, width, height, label):
    right = left + width
    bottom = top + height
    if left < SAFE_L or top < SAFE_T or right > SAFE_R or bottom > SAFE_B:
        raise ValueError(
            f"{label} outside safe area: "
            f"left={left:.2f}, top={top:.2f}, right={right:.2f}, bottom={bottom:.2f}"
        )


def add_bg(slide, dark=False):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = Palette.navy if dark else Palette.paper


def set_run_font(run, size, bold, color, italic=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size,
    bold=False,
    color=Palette.ink,
    align=None,
    italic=False,
    valign=MSO_ANCHOR.TOP,
):
    ensure_box(left, top, width, height, f"text '{text[:24]}'")
    box = slide.shapes.add_textbox(emu(left), emu(top), emu(width), emu(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign

    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if align:
            p.alignment = align
        if p.runs:
            set_run_font(p.runs[0], size, bold, color, italic)
    return box


def add_shape(slide, shape_type, left, top, width, height, fill, line=None, label="shape"):
    ensure_box(left, top, width, height, label)
    shape = slide.shapes.add_shape(shape_type, emu(left), emu(top), emu(width), emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_card(slide, left, top, width, height, fill=Palette.white, line=Palette.line):
    return add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        fill,
        line,
        "card",
    )


def add_title(slide, title, kicker=None, dark=False):
    title_color = Palette.white if dark else Palette.ink
    kicker_color = Palette.muted_dark if dark else Palette.teal_dark
    if kicker:
        add_text(slide, SAFE_L, 0.55, 7.5, 0.28, kicker.upper(), 10, True, kicker_color)
    add_text(slide, SAFE_L, 0.9, 10.9, 0.72, title, 32, True, title_color)
    add_shape(slide, MSO_SHAPE.RECTANGLE, SAFE_L, 1.78, 2.35, 0.06, Palette.teal, label="title line")


def add_footer(slide, number, dark=False):
    color = Palette.muted_dark if dark else Palette.slate
    add_text(slide, SAFE_L, 6.63, 6.2, 0.25, "Frunze Travel · GetVisa · Бишкек", 8.5, False, color)
    add_text(slide, 11.95, 6.63, 0.75, 0.25, f"{number:02}", 8.5, True, color, PP_ALIGN.RIGHT)


def add_bullets(slide, items, left, top, width, height, color=Palette.ink, size=16, accent=Palette.teal):
    gap = height / max(len(items), 1)
    for i, item in enumerate(items):
        y = top + i * gap
        add_shape(slide, MSO_SHAPE.OVAL, left, y + 0.12, 0.12, 0.12, accent, label="bullet")
        add_text(slide, left + 0.28, y, width - 0.28, min(0.45, gap), item, size, False, color)


def set_cell(cell, text, size=12, bold=False, color=Palette.ink, fill=Palette.white, align=PP_ALIGN.CENTER):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.margin_left = Inches(0.05)
    cell.margin_right = Inches(0.05)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)

    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    set_run_font(run, size, bold, color)


def slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, SAFE_L, 0.62, 0.12, 5.95, Palette.teal, label="accent")
    add_text(slide, 0.95, 1.22, 5.4, 0.35, "Frunze Travel", 15, True, Palette.teal)
    add_text(
        slide,
        0.95,
        1.95,
        8.9,
        1.45,
        "Frunze Travel —\nследующий шаг бота",
        34,
        True,
        Palette.white,
    )
    add_text(
        slide,
        0.97,
        3.75,
        8.4,
        0.9,
        "Из потока в 100 сообщений в день — в понятную очередь клиентов с деньгами",
        18,
        False,
        Palette.muted_dark,
    )
    add_card(slide, 9.35, 1.55, 2.95, 3.75, Palette.panel, RGBColor(42, 86, 101))
    add_text(slide, 9.72, 1.95, 2.2, 0.48, "230", 36, True, Palette.white)
    add_text(slide, 9.72, 2.58, 2.2, 0.6, "реальных\nдиалогов", 15, False, Palette.muted_dark)
    add_text(slide, 9.72, 3.65, 2.2, 0.48, "70%", 36, True, Palette.teal)
    add_text(slide, 9.72, 4.28, 2.2, 0.45, "бот провёл сам", 15, False, Palette.muted_dark)
    add_text(slide, 0.95, 6.58, 5.8, 0.25, "Frunze Travel · GetVisa · Бишкек", 9.5, False, Palette.muted_dark)


def slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Что бот уже делает", "реальные данные из WhatsApp")
    metrics = [
        ("230", "диалога\n48–100 новых/день", Palette.teal),
        ("3200", "сообщений\nобработано", Palette.teal_dark),
        ("1400", "ответов бота\nклиентам", Palette.amber),
        ("70%", "диалогов бот\nпровёл сам", Palette.red),
    ]
    for i, (value, label, color) in enumerate(metrics):
        x = SAFE_L + i * 3.08
        add_card(slide, x, 2.25, 2.78, 1.65)
        add_text(slide, x + 0.22, 2.48, 2.25, 0.42, value, 30, True, color)
        add_text(slide, x + 0.22, 3.05, 2.25, 0.52, label, 13, False, Palette.slate)
    add_card(slide, SAFE_L, 4.45, 5.8, 1.25, Palette.green_soft, RGBColor(184, 224, 221))
    add_text(slide, 0.9, 4.72, 5.25, 0.42, "Туры 131 · визы 99", 18, True, Palette.ink)
    add_card(slide, 6.75, 4.45, 5.95, 1.25, Palette.amber_soft, RGBColor(239, 218, 166))
    add_text(slide, 7.05, 4.68, 5.25, 0.54, "24/7: половина обращений вне рабочих часов", 16, True, Palette.ink)
    add_footer(slide, 2)


def slide_3_engagement(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Бот сам делит поток: где деньги, а где шум", "триаж лидов — уже работает")

    steps = [
        ("230", "весь поток за 1-3 июля\nвсе, кто написал", 10.8, Palette.grey_soft, Palette.ink),
        ("160", "бот довёл сам\nбез передачи менеджеру (70%)", 8.65, Palette.teal, Palette.white),
        ("74", "готовы платить\nбот выделил как горячих — им менеджер бегом", 6.35, Palette.panel, Palette.white),
    ]
    y = 2.25
    for i, (number, label, width, fill, color) in enumerate(steps):
        x = SAFE_L + (CONTENT_W - width) / 2
        top = y + i * 1.18
        add_card(slide, x, top, width, 0.92, fill, fill)
        add_text(slide, x + 0.35, top + 0.16, 1.55, 0.42, number, 30, True, color)
        add_text(slide, x + 2.05, top + 0.18, width - 2.45, 0.48, label, 13.5, False, color)

    add_text(
        slide,
        SAFE_L,
        5.95,
        CONTENT_W,
        0.42,
        "Всю разметку бот делает сам, без единой ручной кнопки: 230 — поток, 74 — те, у кого деньги.",
        11,
        False,
        Palette.slate,
        PP_ALIGN.CENTER,
    )
    add_footer(slide, 3)


def slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Откуда брался расход", "почему ~$15 в день — это нормально")
    add_text(
        slide,
        0.6,
        2.05,
        11.4,
        0.45,
        "Деньги уходили не из-за сбоя: бот реально много работал на дорогой модели.",
        18,
        True,
        Palette.ink,
    )
    add_card(slide, 0.6, 2.95, 5.35, 2.95, Palette.grey_soft, Palette.line)
    add_text(slide, 0.95, 3.25, 4.5, 0.36, "Было ~$15/день", 22, True, Palette.ink)
    add_bullets(
        slide,
        [
            "674+ ответов в день",
            "дорогая модель на всём потоке",
            "вся длинная переписка каждый раз",
            "ответы днём и ночью",
        ],
        0.95,
        3.95,
        4.55,
        1.35,
        Palette.slate,
        14.5,
        Palette.grey,
    )
    add_text(slide, 0.95, 5.35, 4.55, 0.28, "Расход = живая работа, не сбой", 14.5, True, Palette.ink)
    add_text(slide, 6.13, 4.05, 0.55, 0.4, "→", 30, True, Palette.teal, PP_ALIGN.CENTER)
    add_card(slide, 6.85, 2.95, 5.85, 2.95, Palette.green_soft, RGBColor(169, 221, 215))
    add_text(slide, 7.2, 3.25, 5.0, 0.36, "Стало ~$1–2,5/день", 22, True, Palette.ink)
    add_bullets(
        slide,
        [
            "дешёвая модель на 80% диалогов",
            "дорогая — только на подборе туров",
            "готовые ответы без AI",
            "сокращение переписки + лимит $10",
        ],
        7.2,
        3.95,
        4.95,
        1.35,
        Palette.slate,
        14.5,
        Palette.teal,
    )
    add_text(slide, 7.2, 5.35, 4.95, 0.28, "Тот же поток — в 2–3 раза дешевле", 14.5, True, Palette.ink)
    add_footer(slide, 4)


def slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Проблема владельца", "деньги тонут в шуме")
    add_text(slide, 0.6, 2.35, 9.9, 0.62, "Из 100 обращений в день реальные деньги — 15–20.", 32, True, Palette.ink)
    add_text(slide, 0.62, 3.22, 8.9, 0.6, "Горячий клиент теряется в общей куче, менеджер видит его слишком поздно.", 17, False, Palette.slate)
    add_card(slide, 0.75, 4.35, 3.35, 1.35, Palette.green_soft, RGBColor(184, 224, 221))
    add_text(slide, 1.0, 4.62, 2.75, 0.34, "15–20", 28, True, Palette.teal)
    add_text(slide, 1.0, 5.08, 2.75, 0.26, "клиентов с деньгами", 13, True, Palette.ink)
    add_card(slide, 4.55, 4.35, 3.35, 1.35, Palette.grey_soft, Palette.line)
    add_text(slide, 4.8, 4.62, 2.75, 0.34, "80+", 28, True, Palette.grey)
    add_text(slide, 4.8, 5.08, 2.75, 0.26, "шум и праздный спрос", 13, True, Palette.ink)
    add_card(slide, 8.45, 4.35, 3.9, 1.35, Palette.white, Palette.line)
    add_text(slide, 8.75, 4.72, 3.2, 0.45, "«Видеть 12 с деньгами, а не читать 100 чатов»", 15, False, Palette.ink, italic=True)
    add_footer(slide, 5)


def slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Решение: светофор", "каждый клиент получает цвет")
    cards = [
        ("Горячий", "даты + бюджет → менеджеру с резюме", Palette.red, Palette.red_soft),
        ("Тёплый", "думает → бот сам напоминает", Palette.amber, Palette.amber_soft),
        ("Просто спросил", "бот отвечает сам, менеджеров не трогает", Palette.grey, Palette.grey_soft),
    ]
    for i, (title, body, accent, fill) in enumerate(cards):
        x = SAFE_L + i * 4.15
        add_card(slide, x, 2.65, 3.75, 3.35, Palette.white, Palette.line)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.18, 2.88, 3.39, 0.16, accent, label="card strip")
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.38, 3.38, 0.68, 0.68, accent, label="status dot")
        add_text(slide, x + 0.38, 4.22, 3.05, 0.36, title, 19, True, Palette.ink)
        add_text(slide, x + 0.38, 4.78, 3.05, 0.75, body, 15, False, Palette.slate)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.18, 5.72, 3.39, 0.08, fill, label="soft strip")
    add_footer(slide, 6)


def slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, True)
    add_title(slide, "Панель владельца «Покупатели сегодня»", "уже работает на проде", dark=True)
    add_card(slide, 0.8, 2.2, 11.45, 3.95, Palette.panel, RGBColor(42, 86, 101))
    add_text(slide, 1.18, 2.52, 6.2, 0.35, "Бот уже разметил поток", 20, True, Palette.white)
    add_text(slide, 1.18, 3.0, 6.6, 0.25, "деньги за 5 секунд, а не 100 чатов — размечено само", 10.5, False, Palette.muted_dark)
    stats = [
        ("74", "готовых\nплатить", Palette.teal),
        ("88", "тёплых\nбот дожимает", Palette.amber),
        ("70", "шум\nотсеян", Palette.grey),
        ("сам", "без ручных\nкнопок", Palette.white),
    ]
    for i, (num, label, color) in enumerate(stats):
        x = 1.18 + i * 2.72
        add_card(slide, x, 3.62, 2.35, 1.55, RGBColor(24, 59, 82), RGBColor(55, 95, 112))
        add_text(slide, x + 0.18, 3.88, 1.82, 0.32, num, 24, True, color)
        add_text(slide, x + 0.18, 4.35, 1.82, 0.42, label, 12.5, False, Palette.white)
    add_text(slide, 1.18, 5.48, 10.6, 0.25, "уже работает: бот сам метит горячих/тёплых/шум и даже исход диалога (оплатил/отказ/пропал)", 10.5, False, Palette.muted_dark)
    add_footer(slide, 7, dark=True)


def slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Что изменится", "результат для бизнеса")
    items = [
        ("Больше оплат", "из того же входящего потока"),
        ("Менеджеры на деньгах", "Адеми, Сезим, Медина, Элиза не перелопачивают шум"),
        ("Репутация цела", "даже простому вопросу отвечают тепло"),
        ("Владелец видит бизнес", "не 100 чатов, а очередь возможностей"),
    ]
    for i, (title, body) in enumerate(items):
        x = SAFE_L + (i % 2) * 6.25
        y = 2.35 + (i // 2) * 2.0
        add_card(slide, x, y, 5.85, 1.45, Palette.white, Palette.line)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.18, y + 0.22, 0.08, 1.0, Palette.teal if i != 1 else Palette.amber, label="mark")
        add_text(slide, x + 0.42, y + 0.25, 5.05, 0.3, title, 18, True, Palette.ink)
        add_text(slide, x + 0.42, y + 0.72, 5.05, 0.42, body, 14.5, False, Palette.slate)
    add_footer(slide, 8)


def slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Как поймём, что работает", "честные метрики")
    metrics = [
        ("Время реакции", "на горячего: минуты, не часы", Palette.red),
        ("Конверсия в оплату", "каждый диалог получает исход", Palette.teal),
        ("Горячих в день", "сколько бот выделил", Palette.amber),
    ]
    for i, (title, body, accent) in enumerate(metrics):
        x = SAFE_L + i * 4.15
        add_card(slide, x, 2.75, 3.75, 2.6, Palette.white, Palette.line)
        add_text(slide, x + 0.35, 3.15, 3.05, 0.42, title, 18, True, Palette.ink)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.35, 3.82, 1.35, 0.08, accent, label="metric line")
        add_text(slide, x + 0.35, 4.2, 3.05, 0.6, body, 15.5, False, Palette.slate)
    add_text(slide, 0.6, 6.05, 10.7, 0.42, "Главный вопрос: стало ли больше оплат из того же потока.", 18, True, Palette.ink)
    add_footer(slide, 9)


def slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, True)
    add_text(slide, 0.6, 1.0, 5.8, 0.28, "Суть в одном предложении", 11, True, Palette.teal)
    add_text(
        slide,
        0.6,
        1.65,
        11.55,
        2.35,
        "Бот становится фильтром: пропускает шум через себя, а вам подаёт людей с деньгами — вовремя и с сохранённой репутацией.",
        32,
        True,
        Palette.white,
    )
    add_card(slide, 0.75, 4.65, 6.4, 1.15, Palette.panel, RGBColor(42, 86, 101))
    add_text(slide, 1.05, 4.96, 5.65, 0.42, "Frunze Travel · GetVisa · Бишкек", 16, True, Palette.white)
    add_text(slide, 0.75, 6.48, 11.0, 0.25, "Данные из рабочей базы за 1–3 июля 2026. Панель «Покупатели сегодня» уже работает на проде.", 8.5, False, Palette.muted_dark)


def slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Сколько людей писали и сколько бот ответил", "реальные данные из прода")

    rows = [
        ("03.07", "82", "581", "487"),
        ("02.07", "100", "981", "776"),
        ("01.07", "48", "244", "137"),
    ]
    headers = ("День", "Новых диалогов", "Сообщений клиентов", "Ответов бота")
    total = ("ВСЕГО", "230 диалогов", "1806", "1400 ответов бота")

    table_left = SAFE_L
    table_top = 2.02
    table_width = CONTENT_W
    table_height = 3.92
    ensure_box(table_left, table_top, table_width, table_height, "activity table")
    frame = slide.shapes.add_table(
        len(rows) + 2,
        len(headers),
        emu(table_left),
        emu(table_top),
        emu(table_width),
        emu(table_height),
    )
    table = frame.table
    col_widths = (1.55, 2.55, 3.75, 4.28)
    for col_idx, width in enumerate(col_widths):
        table.columns[col_idx].width = emu(width)

    for col_idx, header in enumerate(headers):
        set_cell(table.cell(0, col_idx), header, 12, True, Palette.white, Palette.panel)

    for row_idx, row in enumerate(rows, start=1):
        fill = Palette.white if row_idx % 2 else Palette.grey_soft
        for col_idx, value in enumerate(row):
            set_cell(table.cell(row_idx, col_idx), value, 11.5, False, Palette.ink, fill)

    total_row = len(rows) + 1
    for col_idx, value in enumerate(total):
        set_cell(table.cell(total_row, col_idx), value, 12, True, Palette.ink, Palette.amber_soft)

    add_text(
        slide,
        SAFE_L,
        6.08,
        CONTENT_W,
        0.38,
        "Бот в боевом режиме с 1 июля — отвечает 24/7, почти половина обращений приходит вне рабочих часов.",
        10.5,
        False,
        Palette.slate,
    )
    add_footer(slide, 11)


def self_check(prs):
    statuses = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        problems = []
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            left = shape.left / 914400
            top = shape.top / 914400
            width = shape.width / 914400
            height = shape.height / 914400
            right = left + width
            bottom = top + height
            if left < SAFE_L - 0.001 or top < SAFE_T - 0.001 or right > SAFE_R + 0.001 or bottom > SAFE_B + 0.001:
                problems.append(
                    f"shape {shape_idx}: left={left:.2f}, top={top:.2f}, "
                    f"right={right:.2f}, bottom={bottom:.2f}"
                )
        if problems:
            statuses.append((slide_idx, "overflow", problems))
        else:
            statuses.append((slide_idx, "OK", []))
    return statuses


def copy_pptx(src, dst):
    try:
        shutil.copyfile(src, dst)
    except PermissionError:
        src_arg = "'" + str(src).replace("'", "''") + "'"
        dst_arg = "'" + str(dst).replace("'", "''") + "'"
        command = f"Copy-Item -LiteralPath {src_arg} -Destination {dst_arg} -Force -ErrorAction Stop"
        encoded = base64.b64encode(command.encode("utf-16le")).decode("ascii")
        subprocess.run(
            [
                "powershell",
                "-NoProfile",
                "-EncodedCommand",
                encoded,
            ],
            check=True,
            capture_output=True,
            text=True,
        )


def build():
    prs = Presentation()
    prs.slide_width = emu(SLIDE_W)
    prs.slide_height = emu(SLIDE_H)

    for make_slide in (
        slide_1,
        slide_2,
        slide_3_engagement,
        slide_3,
        slide_4,
        slide_5,
        slide_6,
        slide_7,
        slide_8,
        slide_9,
        slide_10,
    ):
        make_slide(prs)

    statuses = self_check(prs)
    for slide_idx, status, problems in statuses:
        print(f"Slide {slide_idx}: {status}")
        for problem in problems:
            print(f"  WARNING {problem}")
    print(f"Slides total: {len(prs.slides)}")

    if any(status != "OK" for _, status, _ in statuses):
        raise SystemExit("Overflow found. Fix coordinates before saving.")

    prs.save(BUILD_PATH)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    FINAL_PATH.parent.mkdir(parents=True, exist_ok=True)
    final_saved = False
    try:
        copy_pptx(BUILD_PATH, FINAL_PATH)
        final_saved = True
    except (PermissionError, subprocess.CalledProcessError) as exc:
        print(f"WARNING: cannot overwrite {FINAL_PATH}: {exc}")
    try:
        copy_pptx(BUILD_PATH, OUT_PATH)
    except (PermissionError, subprocess.CalledProcessError) as exc:
        print(f"WARNING: cannot overwrite {OUT_PATH}: {exc}")
    if final_saved:
        print(f"Saved final: {FINAL_PATH}")
    else:
        print(f"Saved build: {BUILD_PATH}")
    return FINAL_PATH


if __name__ == "__main__":
    build()
