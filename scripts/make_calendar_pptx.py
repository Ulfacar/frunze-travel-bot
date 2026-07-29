"""Сборка PPTX презентации «Календарь звонков» (Frunze Travel).

Концепт Варианта 1: бот = ночной диспетчер, ловит тёплого клиента ночью, бронирует
слот утреннего звонка, к 9:00 отдаёт менеджеру досье. Честно: публичный сайт-календарь
отклонён (разрыв WhatsApp→сайт) → выбор слота В ЧАТЕ; + дыра измеримости продаж.
Переиспользует движок стилей из make_presentation.py (как дека «Дырявое ведро»).

Запуск:  python scripts/make_calendar_pptx.py
Результат: C:\\Users\\alanb\\OneDrive\\Рабочий стол\\Frunze-календарь-звонков.pptx
"""
from __future__ import annotations

import os
import sys
from pathlib import Path

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from make_presentation import (  # noqa: E402
    Palette, SAFE_L, add_bg, add_card, add_shape, add_text,
)

OUT = Path(r"C:\Users\alanb\OneDrive\Рабочий стол\Frunze-календарь-звонков.pptx")
FOOT = "Frunze Travel · Календарь звонков · Бишкек"


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _foot(slide, n, dark=False):
    c = Palette.muted_dark if dark else Palette.slate
    add_text(slide, SAFE_L, 6.66, 7.0, 0.25, FOOT, 8.5, False, c)
    add_text(slide, 11.7, 6.66, 1.0, 0.25, f"{n:02d}", 8.5, True, c, align=PP_ALIGN.RIGHT)


def _head(slide, kicker, title, dark=False):
    kc = Palette.muted_dark if dark else Palette.teal_dark
    tc = Palette.white if dark else Palette.ink
    add_text(slide, SAFE_L, 0.52, 11.5, 0.28, kicker.upper(), 10.5, True, kc)
    add_text(slide, SAFE_L, 0.86, 11.9, 1.0, title, 27, True, tc)
    add_shape(slide, MSO_SHAPE.RECTANGLE, SAFE_L, 1.72, 2.2, 0.055, Palette.teal, label="rule")


def _bubble(slide, x, y, w, h, text, fill, color, italic=False, size=11.5):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, label="bubble")
    add_text(slide, x + 0.14, y + 0.06, w - 0.28, h - 0.12, text, size, False, color,
             italic=italic, valign=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------- slides

def s1_cover(prs):
    s = _blank(prs); add_bg(s, dark=True)
    add_shape(s, MSO_SHAPE.RECTANGLE, SAFE_L, 0.9, 2.2, 0.06, Palette.amber, label="a")
    add_text(s, SAFE_L, 1.15, 11.5, 0.3, "FRUNZE TRAVEL · КАЛЕНДАРЬ ЗВОНКОВ", 11, True, Palette.amber)
    add_text(s, SAFE_L, 1.7, 11.9, 2.0, "Бот ловит ночью.\nЧеловек закрывает утром.",
             40, True, Palette.white)
    add_text(s, SAFE_L, 4.15, 11.0, 1.1,
             "Тур созревает ночью, когда менеджеров нет. Вместо «мы вам перезвоним» бот бронирует "
             "клиенту конкретный слот утреннего звонка — и к 9:00 кладёт менеджеру готовое досье.",
             15, False, Palette.muted_dark)
    for i, t in enumerate(["Ночь поймана", "Утро — стопка звонков", "Показатель show-up"]):
        x = SAFE_L + i * 4.0
        add_shape(s, MSO_SHAPE.OVAL, x, 5.78, 0.12, 0.12, Palette.teal, label="d")
        add_text(s, x + 0.22, 5.66, 3.7, 0.4, t, 12, True, Palette.white)
    _foot(s, 1, dark=True)


def s2_reframe(prs):
    s = _blank(prs); add_bg(s)
    _head(s, "Рефрейм · роль бота", "Бот — ночной диспетчер, а не продавец.")
    # problem card (dark)
    add_card(s, SAFE_L, 1.95, 5.75, 3.9, Palette.navy, Palette.navy)
    add_text(s, SAFE_L + 0.3, 2.12, 5.2, 0.3, "ПОЧЕМУ ТЕЧЁТ", 11, True, Palette.amber)
    for i, t in enumerate([
        "Туры созревают ночью — менеджеры спят",
        "Медиана ответа человека — 3 часа",
        "20% лидов рождаются ночью (22:00–08:00)",
        "Цена — ручное оружие: бот её не называет",
    ]):
        y = 2.6 + i * 0.66
        add_text(s, SAFE_L + 0.3, y, 0.3, 0.4, "—", 13, True, Palette.amber)
        add_text(s, SAFE_L + 0.62, y, 4.9, 0.6, t, 12, False, Palette.white)
    # solution card (light)
    hx = 6.6
    add_card(s, hx, 1.95, 5.75, 3.9, Palette.white, Palette.teal)
    add_text(s, hx + 0.3, 2.12, 5.2, 0.3, "ЧТО ДЕЛАЕТ КАЛЕНДАРЬ", 11, True, Palette.teal_dark)
    for i, t in enumerate([
        "Ночью квалифицирует тёплого клиента",
        "Бронирует слот утреннего звонка",
        "К 9:00 отдаёт менеджеру готовое досье",
        "Утром человек звонит и закрывает",
    ]):
        y = 2.6 + i * 0.66
        add_text(s, hx + 0.3, y, 0.3, 0.4, "→", 13, True, Palette.teal)
        add_text(s, hx + 0.62, y, 4.9, 0.6, t, 12, False, Palette.ink)
    _foot(s, 2)


def s3_roles(prs):
    s = _blank(prs); add_bg(s)
    _head(s, "Как это работает", "Три роли — один конвейер.")
    cols = [
        ("👤 Клиент", "В переписке WhatsApp", Palette.teal,
         "Бот даёт 2–3 свободных времени. Клиент отвечает цифрой — «в 10:00». "
         "Сам выбрал → сам обязался прийти."),
        ("🧑‍💼 Менеджер", "Экран календаря утром", Palette.amber,
         "Приходит в 9:00 — не хаос-инбокс, а стопка назначенных звонков с готовым "
         "досье по каждому: бюджет, даты, кто платит."),
        ("🤖 Бот", "Ночью на автопилоте", Palette.teal_dark,
         "Квалифицирует, бронирует слот наименее загруженному менеджеру, к 9:00 "
         "шлёт «горячий лист» в общий чат."),
    ]
    cw = 3.8
    for i, (name, tag, accent, body) in enumerate(cols):
        x = SAFE_L + i * (cw + 0.36)
        add_card(s, x, 2.0, cw, 3.6, Palette.white, Palette.line)
        add_shape(s, MSO_SHAPE.RECTANGLE, x, 2.0, cw, 0.08, accent, label="top")
        add_text(s, x + 0.25, 2.28, cw - 0.5, 0.45, name, 17, True, Palette.ink)
        add_text(s, x + 0.25, 2.78, cw - 0.5, 0.3, tag.upper(), 9.5, True, accent)
        add_text(s, x + 0.25, 3.25, cw - 0.5, 2.2, body, 12, False, Palette.slate)
    _foot(s, 3)


def s4_client(prs):
    s = _blank(prs); add_bg(s)
    _head(s, "Клиент · выбор в чате", "Слот выбирается прямо в WhatsApp — не на сайте.")
    _bubble(s, SAFE_L, 2.0, 5.6, 0.6, "Анталья, вдвоём, 10–17 августа, до 1500$  · 01:47",
            Palette.grey_soft, Palette.ink)
    _bubble(s, SAFE_L + 0.6, 2.75, 5.6, 0.9,
            "Айгерим, по цене лучше отработает живой менеджер. Адеми свободна завтра утром — "
            "в 10:00 или 11:30? Ответьте 1 или 2  · 01:47", Palette.navy, Palette.white)
    _bubble(s, SAFE_L, 3.85, 2.2, 0.55, "1  · 07:12", Palette.green_soft, Palette.teal_dark)
    _bubble(s, SAFE_L + 0.6, 4.55, 5.6, 0.72,
            "Готово ✅ Адеми позвонит завтра в 10:00. До связи!  · 07:12",
            Palette.navy, Palette.white)
    # right note
    rx = 6.9
    add_card(s, rx, 2.0, 5.6, 1.9, Palette.amber_soft, Palette.amber)
    add_text(s, rx + 0.28, 2.15, 5.1, 0.3, "ПОЧЕМУ НЕ САЙТ", 10.5, True, Palette.amber)
    add_text(s, rx + 0.28, 2.55, 5.1, 1.25,
             "Публичную страницу бронирования на сайте мы отклонили: в Бишкеке клиенты живут "
             "в WhatsApp и не ходят по ссылкам ботов ночью. Каждый уход из чата режет конверсию. "
             "Слот выбирается цифрой — ноль трения.", 12, False, Palette.ink,
             valign=MSO_ANCHOR.MIDDLE)
    add_card(s, rx, 4.1, 5.6, 1.5, Palette.green_soft, Palette.teal)
    add_text(s, rx + 0.28, 4.25, 5.1, 0.3, "ЭФФЕКТ", 10.5, True, Palette.teal_dark)
    add_text(s, rx + 0.28, 4.62, 5.1, 0.9,
             "Кто сам выбрал время — тот приходит. Бронь работает как фильтр: пустышки слот не берут, "
             "менеджер тратит утро только на реальных.", 12, False, Palette.teal_dark,
             valign=MSO_ANCHOR.MIDDLE)
    _foot(s, 4)


def s5_manager(prs):
    s = _blank(prs); add_bg(s)
    _head(s, "Менеджер · утро", "Не 100 «здравствуйте», а стопка назначенных звонков.")
    # appointment stack (left)
    rows = [
        ("10:00", "Айгерим", "Анталья · 1500$ · 10–17 авг", Palette.teal, "✓ подтвердила"),
        ("10:20", "Данияр", "Дубай · бюджет ? · семья 4", Palette.amber, "не ответил"),
        ("10:40", "Салтанат", "Виза США · документы готовы", Palette.teal, "✓ подтвердила"),
        ("11:20", "Нурлан", "Шарм · 2000$ · сентябрь", Palette.teal, "✓ подтвердил"),
    ]
    add_text(s, SAFE_L, 1.95, 6.0, 0.3, "СЕГОДНЯ · АДЕМИ · 4 ЗВОНКА", 10.5, True, Palette.teal_dark)
    for i, (tm, nm, meta, ac, st) in enumerate(rows):
        y = 2.4 + i * 0.85
        add_card(s, SAFE_L, y, 6.0, 0.72, Palette.white, Palette.line)
        add_shape(s, MSO_SHAPE.RECTANGLE, SAFE_L, y, 0.06, 0.72, ac, label="stripe")
        add_text(s, SAFE_L + 0.22, y + 0.16, 0.9, 0.4, tm, 14, True, Palette.ink)
        add_text(s, SAFE_L + 1.25, y + 0.08, 3.0, 0.3, nm, 12.5, True, Palette.ink)
        add_text(s, SAFE_L + 1.25, y + 0.38, 4.4, 0.3, meta, 9.5, False, Palette.slate)
        add_text(s, SAFE_L + 4.4, y + 0.22, 1.5, 0.3, st, 9, True, ac, align=PP_ALIGN.RIGHT)
    # dossier (right)
    rx = 6.9
    add_card(s, rx, 2.4, 5.6, 3.15, Palette.grey_soft, Palette.line)
    add_text(s, rx + 0.28, 2.55, 5.1, 0.3, "ДОСЬЕ К ЗВОНКУ 10:00", 10, True, Palette.teal_dark)
    add_text(s, rx + 0.28, 2.9, 5.1, 0.4, "Айгерим · Анталья", 17, True, Palette.ink)
    for i, (k, v) in enumerate([
        ("Бюджет", "до 1500 $"), ("Даты", "10–17 авг"),
        ("Туристы", "2 взрослых"), ("Кто платит", "она сама"),
    ]):
        y = 3.5 + i * 0.42
        add_text(s, rx + 0.28, y, 2.2, 0.3, k, 11, False, Palette.slate)
        add_text(s, rx + 2.6, y, 2.7, 0.3, v, 11, True, Palette.ink, align=PP_ALIGN.RIGHT)
    add_text(s, rx + 0.28, 5.2, 5.1, 0.3,
             "Заметка бота: сравнивала с 2 агентствами, по нашей цене готова.",
             10, False, Palette.slate, italic=True)
    _foot(s, 5)


def s6_night(prs):
    s = _blank(prs); add_bg(s)
    _head(s, "Ночной бот", "Что бот делает, пока все спят.")
    steps = [
        ("Квалифицирует", "Ведёт диалог, считает готовность: green / warm / шум."),
        ("Предлагает слот", "Дозрел до готового + 4 критерия → даёт 2–3 времени в чате."),
        ("Балансирует", "Слот падает наименее загруженному: Адеми / Сезим / Медина / Элиза."),
        ("Пуш в 9:00", "Собирает горячий лист дня и шлёт менеджерам в Telegram + экран."),
    ]
    for i, (t, d) in enumerate(steps):
        y = 2.05 + i * 1.05
        accent = Palette.amber if i == 3 else Palette.teal
        add_shape(s, MSO_SHAPE.OVAL, SAFE_L, y, 0.5, 0.5, accent, label="n")
        add_text(s, SAFE_L, y + 0.06, 0.5, 0.4, str(i + 1), 15, True, Palette.white, align=PP_ALIGN.CENTER)
        add_text(s, SAFE_L + 0.7, y - 0.02, 6.0, 0.4, t, 14, True, Palette.ink)
        add_text(s, SAFE_L + 0.7, y + 0.38, 6.0, 0.55, d, 11, False, Palette.slate)
    add_card(s, 8.4, 2.05, 4.13, 4.0, Palette.navy, Palette.navy)
    add_text(s, 8.68, 2.25, 3.6, 0.3, "🔥 ГОРЯЧИЙ ЛИСТ · 9:00", 10.5, True, Palette.amber)
    for i, (tm, nm) in enumerate([
        ("10:00", "Айгерим · Анталья"), ("10:20", "Данияр · Дубай"),
        ("10:40", "Салтанат · виза США"), ("11:20", "Нурлан · Шарм"),
    ]):
        y = 2.75 + i * 0.62
        add_text(s, 8.68, y, 0.9, 0.3, tm, 11, True, Palette.teal, align=PP_ALIGN.LEFT)
        add_text(s, 9.55, y, 2.9, 0.3, nm, 11, False, Palette.white)
    add_text(s, 8.68, 5.5, 3.6, 0.5, "→ падает менеджерам в Telegram", 10, False, Palette.muted_dark, italic=True)
    _foot(s, 6)


def s7_kpi(prs):
    s = _blank(prs); add_bg(s)
    _head(s, "Два KPI · разная ответственность", "Бот доводит до слота. Человек — до денег.")
    add_text(s, SAFE_L, 1.85, 11.8, 0.5,
             "Разделение снимает вечный спор «трафик или менеджеры». У каждого свой измеримый результат.",
             13, False, Palette.slate)
    cards = [
        ("KPI БОТА", "% лидов, доведённых до назначенного слота с человеком", Palette.teal, Palette.green_soft),
        ("KPI МЕНЕДЖЕРА", "% слотов, ставших деньгами (show-up → продажа)", Palette.amber, Palette.amber_soft),
    ]
    cw = 5.85
    for i, (who, m, ac, bg) in enumerate(cards):
        x = SAFE_L + i * (cw + 0.4)
        add_card(s, x, 2.6, cw, 2.0, bg, ac)
        add_text(s, x + 0.3, 2.85, cw - 0.6, 0.3, who, 11, True, ac)
        add_text(s, x + 0.3, 3.35, cw - 0.6, 1.0, m, 16, True, Palette.ink, valign=MSO_ANCHOR.MIDDLE)
    add_card(s, SAFE_L, 4.95, 12.13, 1.15, Palette.grey_soft, Palette.line)
    add_text(s, SAFE_L + 0.25, 5.08, 11.6, 0.95,
             "Важно: слот ≠ продажа. Если менеджер не закрывает — точная доставка умножается на ноль. "
             "Поэтому show-up мы обязаны измерять — и здесь вылезает дыра: продажи в системе сейчас "
             "не считаются (см. следующий слайд).", 12, False, Palette.ink, valign=MSO_ANCHOR.MIDDLE)
    _foot(s, 7)


def s8_honest(prs):
    s = _blank(prs); add_bg(s, dark=True)
    add_text(s, SAFE_L, 0.55, 11.0, 0.28, "ЧЕСТНО · ДВА РЕШЕНИЯ", 10.5, True, Palette.amber)
    add_text(s, SAFE_L, 0.9, 11.9, 0.7, "Что мы срезали и что вскрыли.", 27, True, Palette.white)
    add_card(s, SAFE_L, 2.05, 5.9, 4.1, Palette.panel, Palette.panel)
    add_text(s, SAFE_L + 0.3, 2.25, 5.3, 0.3, "СРЕЗАЛИ — САЙТ-КАЛЕНДАРЬ", 10.5, True, Palette.amber)
    add_text(s, SAFE_L + 0.3, 2.65, 5.3, 3.3,
             "Публичная страница бронирования на сайте убила бы конверсию: разрыв WhatsApp→сайт, "
             "паттерн фишинга от бота ночью, риск двойных броней и спама.\n\n"
             "→ Решение: выбор слота цифрой прямо в чате. Меньше кода, меньше трения, тот же результат.",
             13, False, Palette.white)
    add_card(s, 6.73, 2.05, 6.0, 4.1, Palette.panel, Palette.panel)
    add_text(s, 6.73 + 0.3, 2.25, 5.4, 0.3, "ВСКРЫЛИ — ПРОДАЖИ НЕ СЧИТАЮТСЯ", 10.5, True, Palette.teal)
    add_text(s, 6.73 + 0.3, 2.65, 5.4, 3.3,
             "Факт оплаты живёт в Bitrix-«Сделке», которую система не читает. «Продано» размечено "
             "3 из 389 вручную — конверсия недостоверна.\n\n"
             "→ Без чтения сделок назад show-up слотов не с чем сравнить. Это фундамент под весь "
             "календарь — вопрос №8 уже у вас.",
             13, False, Palette.white)
    _foot(s, 8, dark=True)


def s9_phases(prs):
    s = _blank(prs); add_bg(s)
    _head(s, "Фазы · не ждём сложа руки", "Часть уже в проде. Слоты — по вашему сигналу.")
    # phase 0
    add_card(s, SAFE_L, 2.0, 5.9, 3.9, Palette.green_soft, Palette.teal)
    add_text(s, SAFE_L + 0.3, 2.18, 5.3, 0.3, "ФАЗА 0 · УЖЕ РАБОТАЕТ", 11, True, Palette.teal_dark)
    add_text(s, SAFE_L + 0.3, 2.55, 5.3, 0.5, "Утренний «горячий лист»", 17, True, Palette.ink)
    for i, t in enumerate([
        "Бот собирает готовых лидов за ночь",
        "Экран /admin/morning — стопка звонков",
        "Сортировка по ожиданию, чек, «ждёт N»",
        "В проде, на боевых данных сегодня",
    ]):
        y = 3.2 + i * 0.6
        add_text(s, SAFE_L + 0.3, y, 0.3, 0.4, "✓", 13, True, Palette.teal, )
        add_text(s, SAFE_L + 0.62, y, 4.9, 0.55, t, 11.5, False, Palette.ink)
    # phase 1
    add_card(s, 6.73, 2.0, 6.0, 3.9, Palette.white, Palette.amber)
    add_text(s, 6.73 + 0.3, 2.18, 5.4, 0.3, "ФАЗА 1 · ЖДЁТ ГРИШУ + РЕШЕНИЕ ЛПР", 11, True, Palette.amber)
    add_text(s, 6.73 + 0.3, 2.55, 5.4, 0.5, "Слоты бронирования", 17, True, Palette.ink)
    for i, t in enumerate([
        "Таблица слотов + часы менеджеров",
        "Выбор времени в чате WhatsApp",
        "Экран /calendar + напоминание за час",
        "Балансировка нагрузки по менеджерам",
    ]):
        y = 3.2 + i * 0.6
        add_text(s, 6.73 + 0.3, y, 0.3, 0.4, "○", 13, True, Palette.amber)
        add_text(s, 6.73 + 0.62, y, 5.0, 0.55, t, 11.5, False, Palette.ink)
    add_card(s, SAFE_L, 6.05, 12.13, 0.62, Palette.grey_soft, Palette.line)
    add_text(s, SAFE_L + 0.25, 6.13, 11.6, 0.45,
             "~60% фундамента Фазы 1 уже готово (планировщик, мотор готовности, экран-образец). "
             "Реально новое — таблица слотов и напоминание.", 10.5, False, Palette.ink,
             valign=MSO_ANCHOR.MIDDLE)
    _foot(s, 9)


def s10_next(prs):
    s = _blank(prs); add_bg(s, dark=True)
    add_text(s, SAFE_L, 0.55, 11.0, 0.28, "ЧТО НУЖНО ОТ ВАС", 10.5, True, Palette.amber)
    add_text(s, SAFE_L, 0.9, 11.9, 0.7, "Три ответа — и запускаем слоты.", 27, True, Palette.white)
    items = [
        ("1", "Формат календаря", "Свой встроенный + Telegram-пуши (мы за это) или Google Calendar?"),
        ("2", "Доступ к сделкам", "Чтение Bitrix-сделок, чтобы мерить show-up и конверсию (вопрос №8)."),
        ("3", "Telegram-группа", "id чата менеджеров — включим утренний пуш «горячего листа»."),
    ]
    for i, (n, t, d) in enumerate(items):
        y = 2.05 + i * 1.25
        add_card(s, SAFE_L, y, 12.13, 1.05, Palette.panel, Palette.panel)
        add_text(s, SAFE_L + 0.3, y + 0.2, 0.7, 0.65, n, 26, True, Palette.amber)
        add_text(s, SAFE_L + 1.1, y + 0.15, 10.8, 0.4, t, 15, True, Palette.white)
        add_text(s, SAFE_L + 1.1, y + 0.58, 10.8, 0.4, d, 11.5, False, Palette.muted_dark)
    add_text(s, SAFE_L, 5.95, 12.0, 0.6,
             "Фаза 0 уже приносит пользу без этих ответов. Слоты (Фаза 1) — как только решите формат.",
             12, True, Palette.teal, italic=True, align=PP_ALIGN.CENTER)
    _foot(s, 10, dark=True)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for fn in (s1_cover, s2_reframe, s3_roles, s4_client, s5_manager, s6_night,
               s7_kpi, s8_honest, s9_phases, s10_next):
        fn(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"OK: {len(prs.slides._sldIdLst)} слайдов -> {OUT}")


if __name__ == "__main__":
    build()
